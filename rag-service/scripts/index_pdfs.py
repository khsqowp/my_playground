import argparse
import hashlib
import json
import os
import re
import shutil
import subprocess
import tempfile
import uuid
import zlib
from difflib import SequenceMatcher
from pathlib import Path
from typing import Iterable

import fitz
import olefile
import openpyxl
import pytesseract
from docx import Document
from fastembed import SparseTextEmbedding, TextEmbedding
from faster_whisper import WhisperModel
from pptx import Presentation
from qdrant_client import QdrantClient
from qdrant_client.models import (
    Distance, FieldCondition, Filter, MatchValue, PointStruct,
    SparseVector, SparseVectorParams, VectorParams,
)
from tqdm import tqdm

from settings import get_data_root

CACHE_ROOT = Path("/sandbox/cache").resolve()


def data_root() -> Path:
    return get_data_root()


def _emit(event: dict) -> None:
    """Write a single JSON-line progress event to stdout for the web layer to consume."""
    import sys
    print(json.dumps(event, ensure_ascii=False), flush=True, file=sys.stdout)
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".md", ".txt", ".hwp", ".hwpx"}
AUDIO_EXTENSIONS = {".mp3", ".m4a", ".wav", ".flac", ".aac", ".ogg", ".opus"}
VIDEO_EXTENSIONS = {".mp4", ".mov", ".mkv", ".avi", ".webm", ".m4v"}
MEDIA_EXTENSIONS = AUDIO_EXTENSIONS | VIDEO_EXTENSIONS
SUPPORTED_EXTENSIONS = DOCUMENT_EXTENSIONS | MEDIA_EXTENSIONS


def clean_text(text: str) -> str:
    lines = [line.strip() for line in text.splitlines()]
    text = "\n".join(line for line in lines if line)
    return " ".join(text.split())


def chunk_text(text: str, size: int, overlap: int) -> Iterable[str]:
    if size <= overlap:
        raise ValueError("CHUNK_SIZE must be greater than CHUNK_OVERLAP")
    start = 0
    while start < len(text):
        chunk = text[start : start + size].strip()
        if chunk:
            yield chunk
        start += size - overlap


def iter_batches(items: list, batch_size: int) -> Iterable[list]:
    if batch_size <= 0:
        raise ValueError("batch_size must be greater than 0")
    for start in range(0, len(items), batch_size):
        yield items[start : start + batch_size]


def safe_name(value: str) -> str:
    cleaned = re.sub(r"[^0-9A-Za-z가-힣_-]+", "_", value).strip("_")
    return cleaned or "default"


def collection_for_project(project: str) -> str:
    explicit = os.environ.get("COLLECTION_NAME", "").strip()
    default_project = os.environ.get("PROJECT_NAME", "inbox")
    if explicit and project == default_project:
        return explicit
    return f"{safe_name(project)}_docs"


def safe_source_path(path: Path) -> str:
    resolved = path.resolve()
    root = data_root()
    try:
        return str(resolved.relative_to(root))
    except ValueError as exc:
        raise ValueError(f"Refusing to index outside data root {root}: {path}") from exc


def point_id(source: str, locator: str, chunk_idx: int, text: str) -> str:
    raw = f"{source}:{locator}:{chunk_idx}:{text}".encode("utf-8")
    digest = hashlib.sha256(raw).hexdigest()
    return str(uuid.UUID(digest[:32]))


def seconds_to_timestamp(seconds: float) -> str:
    total = max(0, int(seconds))
    hours = total // 3600
    minutes = (total % 3600) // 60
    secs = total % 60
    return f"{hours:02d}:{minutes:02d}:{secs:02d}"


def cache_key(path: Path) -> str:
    source = safe_source_path(path)
    digest = hashlib.sha256(source.encode("utf-8")).hexdigest()[:16]
    return f"{safe_name(path.stem)}-{digest}"


def media_cache_dir(path: Path, project: str) -> Path:
    return CACHE_ROOT / project / "media" / cache_key(path)


def run_command(command: list[str]) -> None:
    result = subprocess.run(command, capture_output=True, text=True, check=False)
    if result.returncode != 0:
        raise RuntimeError(result.stderr.strip() or result.stdout.strip() or f"Command failed: {command}")


def render_page_for_ocr(page: fitz.Page) -> str:
    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
    with tempfile.NamedTemporaryFile(suffix=".png") as image:
        pix.save(image.name)
        return pytesseract.image_to_string(image.name, lang=os.environ.get("OCR_LANG", "kor+eng"))


def extract_pdf(path: Path) -> list[dict]:
    pages = []
    min_text = int(os.environ.get("OCR_MIN_TEXT_CHARS", "1"))
    with fitz.open(path) as doc:
        for page_idx, page in enumerate(doc, start=1):
            text = clean_text(page.get_text("text"))
            ocr_used = False
            if len(text) < min_text and os.environ.get("ENABLE_OCR", "1") == "1":
                text = clean_text(render_page_for_ocr(page))
                ocr_used = True
            if text:
                pages.append({"locator": f"page {page_idx}", "page": page_idx, "text": text, "ocr": ocr_used})
    return pages


def extract_docx(path: Path) -> list[dict]:
    doc = Document(path)
    parts = [paragraph.text for paragraph in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return [{"locator": "document", "page": None, "text": clean_text("\n".join(parts)), "ocr": False}]


def extract_pptx(path: Path) -> list[dict]:
    prs = Presentation(path)
    pages = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        text = clean_text("\n".join(parts))
        if text:
            pages.append({"locator": f"slide {idx}", "page": idx, "text": text, "ocr": False})
    return pages


def extract_xlsx(path: Path) -> list[dict]:
    workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    pages = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None]
            if values:
                rows.append(" | ".join(values))
        text = clean_text("\n".join(rows))
        if text:
            pages.append({"locator": f"sheet {sheet.title}", "page": None, "text": text, "ocr": False})
    workbook.close()
    return pages


def extract_text(path: Path) -> list[dict]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return [{"locator": "document", "page": None, "text": clean_text(text), "ocr": False}]


def extract_hwp(path: Path) -> list[dict]:
    if path.suffix.lower() == ".hwpx":
        return extract_hwpx_with_libreoffice(path)
    if not olefile.isOleFile(path):
        return []
    ole = olefile.OleFileIO(path)
    compressed = False
    if ole.exists("FileHeader"):
        header = ole.openstream("FileHeader").read()
        compressed = bool(header[36] & 1)
    body_streams = sorted(
        entry for entry in ole.listdir(streams=True, storages=False) if len(entry) == 2 and entry[0] == "BodyText"
    )
    texts = []
    for entry in body_streams:
        data = ole.openstream(entry).read()
        if compressed:
            data = zlib.decompress(data, -15)
        texts.append(data.decode("utf-16le", errors="ignore"))
    return [{"locator": "document", "page": None, "text": clean_text("\n".join(texts)), "ocr": False}]


def media_duration_seconds(path: Path) -> float:
    result = subprocess.run(
        [
            "ffprobe",
            "-v",
            "error",
            "-show_entries",
            "format=duration",
            "-of",
            "default=noprint_wrappers=1:nokey=1",
            str(path),
        ],
        capture_output=True,
        text=True,
        check=False,
    )
    if result.returncode != 0:
        raise RuntimeError(result.stderr.strip() or f"ffprobe failed for {path}")
    return float(result.stdout.strip())


def extract_audio_chunk(source: Path, output: Path, start_seconds: int, chunk_seconds: int) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    if output.exists() and output.stat().st_size > 0:
        return
    run_command(
        [
            "ffmpeg",
            "-hide_banner",
            "-loglevel",
            "error",
            "-y",
            "-ss",
            str(start_seconds),
            "-t",
            str(chunk_seconds),
            "-i",
            str(source),
            "-vn",
            "-ac",
            "1",
            "-ar",
            "16000",
            "-c:a",
            "pcm_s16le",
            str(output),
        ]
    )


def transcribe_audio_chunk(
    model: WhisperModel, audio_path: Path, start_offset: int, context_prompt: str | None = None
) -> list[dict]:
    language = os.environ.get("STT_LANGUAGE", "ko").strip() or None
    beam_size = int(os.environ.get("STT_BEAM_SIZE", "10"))
    min_avg_logprob = float(os.environ.get("STT_MIN_AVG_LOGPROB", "-0.8"))
    max_no_speech_prob = float(os.environ.get("STT_MAX_NO_SPEECH_PROB", "0.3"))
    temperature = float(os.environ.get("STT_TEMPERATURE", "0.0"))
    compression_ratio_threshold = float(os.environ.get("STT_COMPRESSION_RATIO_THRESHOLD", "2.2"))
    default_prompt = os.environ.get("STT_INITIAL_PROMPT", "다음은 한국어 음성입니다.").strip() or None
    effective_prompt = context_prompt if context_prompt is not None else default_prompt
    segments, info = model.transcribe(
        str(audio_path),
        language=language,
        beam_size=beam_size,
        vad_filter=True,
        word_timestamps=False,
        condition_on_previous_text=True,
        temperature=temperature,
        compression_ratio_threshold=compression_ratio_threshold,
        log_prob_threshold=min_avg_logprob,
        no_speech_threshold=max_no_speech_prob,
        initial_prompt=effective_prompt,
    )
    records = []
    for segment in segments:
        avg_logprob = getattr(segment, "avg_logprob", None)
        no_speech_prob = getattr(segment, "no_speech_prob", None)
        if avg_logprob is not None and avg_logprob < min_avg_logprob:
            continue
        if no_speech_prob is not None and no_speech_prob > max_no_speech_prob:
            continue
        text = clean_text(segment.text)
        if not text:
            continue
        start = start_offset + float(segment.start)
        end = start_offset + float(segment.end)
        records.append(
            {
                "start": start,
                "end": end,
                "text": text,
                "avg_logprob": avg_logprob,
                "no_speech_prob": no_speech_prob,
                "language": getattr(info, "language", language),
            }
        )
    return records


# ── Korean spacing correction ─────────────────────────────────────────────────

_kiwi_instance = None


def _get_kiwi():
    global _kiwi_instance
    if _kiwi_instance is None and os.environ.get("STT_ENABLE_KOSPACING", "1") == "1":
        try:
            from kiwipiepy import Kiwi
            _kiwi_instance = Kiwi()
        except ImportError:
            pass
    return _kiwi_instance


def apply_korean_spacing(text: str) -> str:
    if len(text) < 4:
        return text
    kiwi = _get_kiwi()
    if kiwi is None:
        return text
    try:
        return kiwi.space(text)
    except Exception:
        return text


# ── STT ensemble helpers ──────────────────────────────────────────────────────

def get_stt_model_names() -> list[str]:
    models_str = os.environ.get("STT_MODELS", "").strip()
    if models_str:
        return [m.strip() for m in models_str.split(",") if m.strip()]
    return [os.environ.get("STT_MODEL", "small")]


def _model_slug(model_name: str) -> str:
    return re.sub(r"[^0-9A-Za-z]", "_", model_name)


def model_raw_cache_paths(cache_dir: Path, model_name: str) -> tuple[Path, Path]:
    slug = _model_slug(model_name)
    return cache_dir / f"raw_{slug}_manifest.json", cache_dir / f"raw_{slug}.json"


def per_model_manifest(path: Path, model_name: str) -> dict:
    stat = path.stat()
    return {
        "source": safe_source_path(path),
        "size": stat.st_size,
        "mtime_ns": stat.st_mtime_ns,
        "model": model_name,
        "device": os.environ.get("STT_DEVICE", "cpu"),
        "compute_type": os.environ.get("STT_COMPUTE_TYPE", "int8"),
        "language": os.environ.get("STT_LANGUAGE", "ko"),
        "chunk_seconds": int(os.environ.get("STT_CHUNK_SECONDS", "180")),
        "beam_size": int(os.environ.get("STT_BEAM_SIZE", "10")),
        "min_avg_logprob": float(os.environ.get("STT_MIN_AVG_LOGPROB", "-0.8")),
        "max_no_speech_prob": float(os.environ.get("STT_MAX_NO_SPEECH_PROB", "0.3")),
        "temperature": float(os.environ.get("STT_TEMPERATURE", "0.0")),
        "compression_ratio_threshold": float(os.environ.get("STT_COMPRESSION_RATIO_THRESHOLD", "2.2")),
        "initial_prompt": os.environ.get("STT_INITIAL_PROMPT", "다음은 한국어 음성입니다.").strip(),
    }


def load_per_model_cache(cache_dir: Path, model_name: str, manifest: dict) -> list[dict] | None:
    manifest_path, data_path = model_raw_cache_paths(cache_dir, model_name)
    if not manifest_path.exists() or not data_path.exists():
        return None
    if json.loads(manifest_path.read_text(encoding="utf-8")) != manifest:
        return None
    return json.loads(data_path.read_text(encoding="utf-8"))


def save_per_model_cache(cache_dir: Path, model_name: str, manifest: dict, records: list[dict]) -> None:
    cache_dir.mkdir(parents=True, exist_ok=True)
    manifest_path, data_path = model_raw_cache_paths(cache_dir, model_name)
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    data_path.write_text(json.dumps(records, ensure_ascii=False, indent=2), encoding="utf-8")


def ensemble_manifest(path: Path, model_names: list[str]) -> dict:
    stat = path.stat()
    return {
        "source": safe_source_path(path),
        "size": stat.st_size,
        "mtime_ns": stat.st_mtime_ns,
        "models": sorted(model_names),
        "device": os.environ.get("STT_DEVICE", "cpu"),
        "compute_type": os.environ.get("STT_COMPUTE_TYPE", "int8"),
        "language": os.environ.get("STT_LANGUAGE", "ko"),
        "chunk_seconds": int(os.environ.get("STT_CHUNK_SECONDS", "180")),
        "beam_size": int(os.environ.get("STT_BEAM_SIZE", "10")),
        "min_avg_logprob": float(os.environ.get("STT_MIN_AVG_LOGPROB", "-0.8")),
        "max_no_speech_prob": float(os.environ.get("STT_MAX_NO_SPEECH_PROB", "0.3")),
        "temperature": float(os.environ.get("STT_TEMPERATURE", "0.0")),
        "compression_ratio_threshold": float(os.environ.get("STT_COMPRESSION_RATIO_THRESHOLD", "2.2")),
        "initial_prompt": os.environ.get("STT_INITIAL_PROMPT", "다음은 한국어 음성입니다.").strip(),
        "min_votes": int(os.environ.get("STT_MIN_VOTES", "2")),
        "similarity_threshold": float(os.environ.get("STT_SIMILARITY_THRESHOLD", "0.6")),
        "align_tolerance": float(os.environ.get("STT_ALIGN_TOLERANCE", "0.5")),
    }


def load_cached_transcript(cache_dir: Path, manifest: dict) -> list[dict] | None:
    manifest_path = cache_dir / "stt_manifest.json"
    transcript_path = cache_dir / "transcript.json"
    if not manifest_path.exists() or not transcript_path.exists():
        return None
    if json.loads(manifest_path.read_text(encoding="utf-8")) != manifest:
        return None
    return json.loads(transcript_path.read_text(encoding="utf-8"))


def save_transcript_cache(cache_dir: Path, manifest: dict, records: list[dict]) -> None:
    cache_dir.mkdir(parents=True, exist_ok=True)
    (cache_dir / "stt_manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    (cache_dir / "transcript.json").write_text(json.dumps(records, ensure_ascii=False, indent=2), encoding="utf-8")
    lines = []
    for record in records:
        confidence = record.get("stt_confidence", "")
        votes = record.get("votes", "")
        tag = f"[{confidence}/{votes}] " if confidence else ""
        lines.append(f"[{seconds_to_timestamp(record['start'])}-{seconds_to_timestamp(record['end'])}] {tag}{record['text']}")
    (cache_dir / "transcript.txt").write_text("\n".join(lines), encoding="utf-8")


# ── Voting / ensemble logic ───────────────────────────────────────────────────

def _text_similarity(a: str, b: str) -> float:
    return SequenceMatcher(None, a, b).ratio()


def _confidence_tier(votes: int, total: int) -> str:
    ratio = votes / max(total, 1)
    if ratio == 1.0:
        return "CONFIRMED"
    if ratio >= 0.75:
        return "HIGH"
    if ratio >= 0.5:
        return "MEDIUM"
    return "LOW"


def _merge_into_events(segments: list[dict], tolerance: float) -> list[list[dict]]:
    if not segments:
        return []
    sorted_segs = sorted(segments, key=lambda s: s["start"])
    events, current = [], [sorted_segs[0]]
    window_end = sorted_segs[0]["end"]
    for seg in sorted_segs[1:]:
        if seg["start"] <= window_end + tolerance:
            current.append(seg)
            window_end = max(window_end, seg["end"])
        else:
            events.append(current)
            current = [seg]
            window_end = seg["end"]
    events.append(current)
    return events


def _cluster_by_similarity(segs: list[dict], threshold: float) -> list[list[dict]]:
    clusters: list[list[dict]] = []
    assigned = [False] * len(segs)
    for i, seg_i in enumerate(segs):
        if assigned[i]:
            continue
        cluster = [seg_i]
        assigned[i] = True
        for j, seg_j in enumerate(segs):
            if assigned[j]:
                continue
            if _text_similarity(seg_i["text"], seg_j["text"]) >= threshold:
                cluster.append(seg_j)
                assigned[j] = True
        clusters.append(cluster)
    return clusters


def vote_and_merge(
    model_segments: dict[str, list[dict]],
    min_votes: int,
    similarity_threshold: float,
    align_tolerance: float,
) -> list[dict]:
    n_models = len(model_segments)
    effective_min = min(min_votes, n_models)
    all_segs: list[dict] = []
    for model_name, segs in model_segments.items():
        for seg in segs:
            all_segs.append({**seg, "_model": model_name})
    if not all_segs:
        return []
    events = _merge_into_events(all_segs, align_tolerance)
    accepted: list[dict] = []
    for event in events:
        clusters = _cluster_by_similarity(event, similarity_threshold)
        best_cluster = max(clusters, key=lambda c: len(set(s["_model"] for s in c)))
        distinct_models = len(set(s["_model"] for s in best_cluster))
        if distinct_models < effective_min:
            continue
        best_seg = max(best_cluster, key=lambda s: (s.get("avg_logprob") or -999.0))
        accepted.append({
            "start": best_seg["start"],
            "end": best_seg["end"],
            "text": best_seg["text"],
            "avg_logprob": best_seg.get("avg_logprob"),
            "no_speech_prob": best_seg.get("no_speech_prob"),
            "language": best_seg.get("language"),
            "votes": distinct_models,
            "total_models": n_models,
            "stt_confidence": _confidence_tier(distinct_models, n_models),
            "model": best_seg["_model"],
        })
    accepted.sort(key=lambda s: s["start"])
    return accepted


# ── Per-model transcription with cache ───────────────────────────────────────

def _transcribe_full(model_name: str, path: Path, cache_dir: Path, duration: float) -> list[dict]:
    manifest = per_model_manifest(path, model_name)
    cached = load_per_model_cache(cache_dir, model_name, manifest)
    if cached is not None:
        print(f"  [{model_name}] loaded from cache ({len(cached)} segments)")
        return cached
    chunk_seconds = int(os.environ.get("STT_CHUNK_SECONDS", "300"))
    print(f"  [{model_name}] loading model...")
    model = WhisperModel(
        model_name,
        device=os.environ.get("STT_DEVICE", "cpu"),
        compute_type=os.environ.get("STT_COMPUTE_TYPE", "int8"),
        download_root=str(CACHE_ROOT / "whisper"),
    )
    records: list[dict] = []
    context_prompt: str | None = None
    for start in range(0, max(1, int(duration)), chunk_seconds):
        chunk_path = cache_dir / "chunks" / f"{start:08d}.wav"
        extract_audio_chunk(path, chunk_path, start, chunk_seconds)
        chunk_records = transcribe_audio_chunk(model, chunk_path, start, context_prompt=context_prompt)
        records.extend(chunk_records)
        if chunk_records:
            last_texts = " ".join(r["text"] for r in chunk_records[-3:])
            context_prompt = last_texts[-200:]
    del model
    save_per_model_cache(cache_dir, model_name, manifest, records)
    print(f"  [{model_name}] transcribed {len(records)} segments")
    return records


# ── Media extraction entry point ─────────────────────────────────────────────

def extract_media(path: Path, project: str) -> list[dict]:
    if os.environ.get("ENABLE_STT", "1") != "1":
        return []
    cache_dir = media_cache_dir(path, project)
    model_names = get_stt_model_names()
    e_manifest = ensemble_manifest(path, model_names)
    cached = load_cached_transcript(cache_dir, e_manifest)
    if cached is not None:
        records = cached
    else:
        duration = media_duration_seconds(path)
        model_segments: dict[str, list[dict]] = {}
        for model_name in model_names:
            print(f"[STT] {path.name} → {model_name}")
            model_segments[model_name] = _transcribe_full(model_name, path, cache_dir, duration)

        if len(model_names) == 1:
            records = model_segments[model_names[0]]
            for rec in records:
                rec.setdefault("votes", 1)
                rec.setdefault("total_models", 1)
                rec.setdefault("stt_confidence", "CONFIRMED")
                rec.setdefault("model", model_names[0])
        else:
            min_votes = int(os.environ.get("STT_MIN_VOTES", "2"))
            similarity_threshold = float(os.environ.get("STT_SIMILARITY_THRESHOLD", "0.6"))
            align_tolerance = float(os.environ.get("STT_ALIGN_TOLERANCE", "0.5"))
            records = vote_and_merge(model_segments, min_votes, similarity_threshold, align_tolerance)
            total_raw = sum(len(s) for s in model_segments.values())
            print(
                f"[STT] ensemble: {total_raw} raw segments across {len(model_names)} models → "
                f"{len(records)} accepted after voting (min_votes={min_votes})"
            )

        save_transcript_cache(cache_dir, e_manifest, records)

    pages = []
    for idx, record in enumerate(records, start=1):
        start = seconds_to_timestamp(record["start"])
        end = seconds_to_timestamp(record["end"])
        pages.append(
            {
                "locator": f"{start}-{end}",
                "page": None,
                "text": apply_korean_spacing(record["text"]),
                "ocr": False,
                "start_seconds": record["start"],
                "end_seconds": record["end"],
                "stt": True,
                "stt_model": record.get("model", model_names[0]),
                "stt_confidence": record.get("stt_confidence"),
                "stt_votes": record.get("votes"),
            }
        )
    return pages


def extract_hwpx_with_libreoffice(path: Path) -> list[dict]:
    with tempfile.TemporaryDirectory() as temp_dir:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "txt:Text", "--outdir", temp_dir, str(path)],
            capture_output=True,
            text=True,
            check=False,
        )
        txt_files = list(Path(temp_dir).glob("*.txt"))
        if result.returncode != 0 or not txt_files:
            return []
        return [{"locator": "document", "page": None, "text": clean_text(txt_files[0].read_text(errors="ignore")), "ocr": False}]


def extract_records(path: Path, project: str) -> list[dict]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".xlsx":
        return extract_xlsx(path)
    if suffix in {".md", ".txt"}:
        return extract_text(path)
    if suffix in {".hwp", ".hwpx"}:
        return extract_hwp(path)
    if suffix in MEDIA_EXTENSIONS:
        return extract_media(path, project)
    return []


def extract_chunks(path: Path, chunk_size: int, overlap: int, min_chunk_chars: int, project: str) -> list[dict]:
    source = safe_source_path(path)
    chunks: list[dict] = []
    for record in extract_records(path, project):
        text = record["text"]
        for chunk_idx, chunk in enumerate(chunk_text(text, chunk_size, overlap)):
            if len(chunk) < min_chunk_chars:
                continue
            chunks.append(
                {
                    "project": project,
                    "source": source,
                    "locator": record["locator"],
                    "page": record["page"],
                    "chunk_idx": chunk_idx,
                    "extension": path.suffix.lower(),
                    "ocr": record["ocr"],
                    "stt": record.get("stt", False),
                    "stt_model": record.get("stt_model"),
                    "stt_confidence": record.get("stt_confidence"),
                    "stt_votes": record.get("stt_votes"),
                    "start_seconds": record.get("start_seconds"),
                    "end_seconds": record.get("end_seconds"),
                    "text": chunk,
                }
            )
    return chunks


def iter_supported_files(project_dir: Path) -> list[Path]:
    files: list[Path] = []
    for dirpath, dirnames, filenames in os.walk(project_dir, onerror=lambda exc: _emit({
        "event": "scan_error",
        "path": getattr(exc, "filename", ""),
        "error": str(exc),
    })):
        dirnames[:] = [name for name in dirnames if not name.startswith(".")]
        for filename in filenames:
            path = Path(dirpath) / filename
            if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
                continue
            try:
                if path.is_file():
                    files.append(path)
            except OSError as exc:
                _emit({"event": "scan_error", "path": str(path), "error": str(exc)})
    return sorted(files)


def file_fingerprint(path: Path) -> dict:
    stat = path.stat()
    return {
        "path": safe_source_path(path),
        "size": stat.st_size,
        "mtime_ns": stat.st_mtime_ns,
        "extension": path.suffix.lower(),
    }


def manifest_path(project: str) -> Path:
    return CACHE_ROOT / project / "index_manifest.json"


def load_manifest(project: str) -> dict:
    path = manifest_path(project)
    if not path.exists():
        return {"files": {}}
    data = json.loads(path.read_text(encoding="utf-8"))
    files = data.get("files", {})
    if isinstance(files, list):
        return {"files": {item["path"]: item for item in files if isinstance(item, dict) and "path" in item}}
    if not isinstance(files, dict):
        return {"files": {}}
    return {"files": files}


def save_manifest(project: str, manifest: dict) -> None:
    path = manifest_path(project)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")


def ensure_collection(client: QdrantClient, collection: str, vector_size: int, recreate: bool) -> bool:
    """Create or validate the collection. Returns True if a full reindex is required (schema migration)."""
    if recreate and client.collection_exists(collection):
        client.delete_collection(collection)
    if not client.collection_exists(collection):
        client.create_collection(
            collection_name=collection,
            vectors_config={"dense": VectorParams(size=vector_size, distance=Distance.COSINE)},
            sparse_vectors_config={"sparse": SparseVectorParams()},
        )
        return False
    info = client.get_collection(collection)
    if not isinstance(info.config.params.vectors, dict):
        client.delete_collection(collection)
        client.create_collection(
            collection_name=collection,
            vectors_config={"dense": VectorParams(size=vector_size, distance=Distance.COSINE)},
            sparse_vectors_config={"sparse": SparseVectorParams()},
        )
        return True
    return False


def delete_source(client: QdrantClient, collection: str, source: str) -> None:
    if not client.collection_exists(collection):
        return
    client.delete(
        collection_name=collection,
        points_selector=Filter(must=[FieldCondition(key="source", match=MatchValue(value=source))]),
    )


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--project", default="inbox")
    parser.add_argument("--recreate", action="store_true")
    parser.add_argument("--skip-unchanged", action="store_true")
    args = parser.parse_args()

    root = data_root()
    project_dir = (root / args.project).resolve()
    try:
        project_dir.relative_to(root)
    except ValueError:
        raise SystemExit(f"Project directory not found under data root {root}: {args.project}")
    if not project_dir.is_dir():
        raise SystemExit(f"Project directory not found under data root {root}: {args.project}")

    qdrant_url = os.environ["QDRANT_URL"]
    collection = collection_for_project(args.project)
    model_name = os.environ["EMBEDDING_MODEL"]
    chunk_size = int(os.environ.get("CHUNK_SIZE", "1400"))
    overlap = int(os.environ.get("CHUNK_OVERLAP", "200"))
    min_chunk_chars = int(os.environ.get("MIN_CHUNK_CHARS", "80"))
    index_batch_size = int(os.environ.get("INDEX_BATCH_SIZE", "128"))

    files = iter_supported_files(project_dir)
    if not files:
        _emit({"event": "no_files", "project": args.project})
        return

    previous = load_manifest(args.project)
    previous_files = previous.get("files", {})
    current_files = {}
    readable_files = []
    for path in files:
        try:
            source = safe_source_path(path)
            current_files[source] = file_fingerprint(path)
            readable_files.append(path)
        except OSError as exc:
            _emit({"event": "scan_error", "path": str(path), "error": str(exc)})
    files = readable_files
    changed = [path for path in files if args.recreate or previous_files.get(safe_source_path(path)) != current_files[safe_source_path(path)]]
    deleted = sorted(set(previous_files) - set(current_files))

    # Check collection schema before the skip-unchanged early exit so that a
    # schema migration (old → hybrid named-vector format) is never skipped.
    qdrant_client_early = QdrantClient(url=qdrant_url)
    schema_needs_migration = (
        qdrant_client_early.collection_exists(collection)
        and not isinstance(qdrant_client_early.get_collection(collection).config.params.vectors, dict)
    )

    if args.skip_unchanged and not changed and not deleted and not schema_needs_migration:
        _emit({"event": "no_changes", "project": args.project})
        return

    embedder = TextEmbedding(model_name=model_name)
    sparse_embedder = SparseTextEmbedding(model_name="Qdrant/bm25")
    probe = next(embedder.embed(["dimension probe"]))
    vector_size = len(probe)

    client = qdrant_client_early
    needs_full_reindex = ensure_collection(client, collection, vector_size, args.recreate)
    if needs_full_reindex:
        _emit({"event": "migration", "message": "스키마 마이그레이션 (→ hybrid), 전체 재색인 시작"})
        changed = files
        deleted = []

    _emit({
        "event": "start",
        "total_files": len(changed),
        "total_deleted": len(deleted),
        "project": args.project,
    })

    for i, source in enumerate(deleted, 1):
        _emit({"event": "delete", "index": i, "total": len(deleted), "source": source})
        delete_source(client, collection, source)

    total_chunks = 0
    files_done = 0
    files_failed = 0
    files_skipped = 0
    total_files = len(changed)

    for file_index, path in enumerate(changed, 1):
        source = safe_source_path(path)
        filename = path.name
        _emit({
            "event": "file_start",
            "file_index": file_index,
            "total_files": total_files,
            "filename": filename,
            "source": source,
        })
        delete_source(client, collection, source)
        try:
            chunks = extract_chunks(path, chunk_size, overlap, min_chunk_chars, args.project)
        except Exception as exc:
            files_failed += 1
            _emit({
                "event": "file_error",
                "file_index": file_index,
                "total_files": total_files,
                "filename": filename,
                "source": source,
                "error": str(exc),
            })
            continue
        if not chunks:
            files_skipped += 1
            _emit({
                "event": "file_skip",
                "file_index": file_index,
                "total_files": total_files,
                "filename": filename,
                "source": source,
                "reason": "텍스트 추출 결과 없음",
            })
            continue
        file_points = 0
        for chunk_batch in iter_batches(chunks, index_batch_size):
            texts = [chunk["text"] for chunk in chunk_batch]
            vectors = list(tqdm(embedder.embed(texts), total=len(texts), desc=filename, file=__import__("sys").stderr))
            sparse_vectors = list(sparse_embedder.embed(texts))
            points = [
                PointStruct(
                    id=point_id(chunk["source"], chunk["locator"], chunk["chunk_idx"], chunk["text"]),
                    vector={
                        "dense": vector.tolist(),
                        "sparse": SparseVector(
                            indices=sv.indices.tolist(),
                            values=sv.values.tolist(),
                        ),
                    },
                    payload=chunk,
                )
                for chunk, vector, sv in zip(chunk_batch, vectors, sparse_vectors)
            ]
            client.upsert(collection_name=collection, points=points)
            file_points += len(points)
        total_chunks += file_points
        files_done += 1
        _emit({
            "event": "file_done",
            "file_index": file_index,
            "total_files": total_files,
            "filename": filename,
            "source": source,
            "chunks": file_points,
        })

    save_manifest(args.project, {"files": current_files})
    _emit({
        "event": "done",
        "project": args.project,
        "files_done": files_done,
        "files_failed": files_failed,
        "files_skipped": files_skipped,
        "files_deleted": len(deleted),
        "total_chunks": total_chunks,
        "ok": files_failed == 0,
    })


if __name__ == "__main__":
    main()
